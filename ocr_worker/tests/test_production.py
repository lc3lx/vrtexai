"""Integration, E2E, fault-injection, template, review, and stress tests."""
from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from openpyxl import Workbook, load_workbook
from PIL import Image, ImageDraw, ImageFont

from classify import classify_document
from invoice import parse_invoice_table
from table_detect import cluster_words_to_table, detect_cells
from templates import app_templates_path, load_templates, match_template, save_templates


def _font(size: int = 24):
    for candidate in (
        Path(r"C:\Windows\Fonts\arial.ttf"),
        Path(r"C:\Windows\Fonts\tahoma.ttf"),
        Path(r"C:\Windows\Fonts\seguiui.ttf"),
    ):
        if candidate.is_file():
            return ImageFont.truetype(str(candidate), size)
    return ImageFont.load_default()


def run_worker(inputs: list[Path], output_dir: Path, master: Path | None = None, workers: int = 2) -> dict:
    output_dir.mkdir(parents=True, exist_ok=True)
    request_path = output_dir.parent / f"req-{output_dir.name}.json"
    result_path = output_dir.parent / f"res-{output_dir.name}.json"
    payload = {
        "inputPaths": [str(path) for path in inputs],
        "outputDirectory": str(output_dir),
        "resultPath": str(result_path),
        "maxWorkers": workers,
    }
    if master:
        payload["masterDataDirectory"] = str(master)
    request_path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    env = os.environ.copy()
    env["PYTHONPATH"] = str(ROOT)
    tess = ROOT.parent / "runtime" / "tesseract" / "tesseract.exe"
    tessdata = ROOT.parent / "runtime" / "tesseract" / "tessdata"
    if tess.is_file():
        env["TESSERACT_CMD"] = str(tess)
        env["TESSDATA_PREFIX"] = str(tessdata)
    completed = subprocess.run(
        [sys.executable, str(ROOT / "main.py"), "--request", str(request_path)],
        capture_output=True, text=True, env=env, timeout=240,
    )
    if completed.returncode != 0:
        raise AssertionError(f"worker failed: {completed.stderr[-1500:] or completed.stdout[-1500:]}")
    return json.loads(result_path.read_text(encoding="utf-8-sig"))


class IntakeAndBatchTests(unittest.TestCase):
    def setUp(self):
        self.temp = Path(tempfile.mkdtemp(prefix="vertex-prod-"))
        self.output = self.temp / "out"

    def tearDown(self):
        shutil.rmtree(self.temp, ignore_errors=True)

    def _write_csv(self, name: str, rows: list[list[str]]) -> Path:
        path = self.temp / name
        path.write_text("\n".join(",".join(row) for row in rows), encoding="utf-8")
        return path

    def test_unsupported_does_not_stop_batch(self):
        good = self._write_csv("ok.csv", [["name", "qty"], ["A", "1"], ["B", "2"]])
        bad = self.temp / "notes.xyz"
        bad.write_text("not a document", encoding="utf-8")
        result = run_worker([good, bad], self.output, workers=1)
        self.assertEqual(result["filesProcessed"], 2)
        self.assertEqual(result["filesFailed"], 1)
        self.assertGreaterEqual(result["recordsExtracted"], 2)
        status = load_workbook(result["statusReportPath"], read_only=True, data_only=True)
        labels = [row[1] for row in status.active.iter_rows(min_row=2, values_only=True)]
        status.close()
        self.assertIn("unsupported", labels)

    def test_corrupt_pdf_and_image_do_not_crash(self):
        pdf = self.temp / "broken.pdf"
        pdf.write_bytes(b"%PDF-1.1 broken")
        img = self.temp / "broken.png"
        img.write_bytes(b"\x89PNG not really")
        good = self._write_csv("ok.csv", [["a", "b"], ["1", "2"]])
        result = run_worker([pdf, img, good], self.output, workers=1)
        self.assertEqual(result["filesProcessed"], 3)
        self.assertGreaterEqual(result["filesFailed"], 1)
        self.assertGreaterEqual(result["recordsExtracted"], 1)

    def test_corrupt_templates_do_not_crash(self):
        # Keep this fault-injection test isolated from the user's real
        # templates folder (which can also be locked by a running app).
        previous_local_app_data = os.environ.get("LOCALAPPDATA")
        os.environ["LOCALAPPDATA"] = str(self.temp / "localappdata")
        try:
            app = app_templates_path()
            backup = None
            if app.is_file():
                backup = app.read_text(encoding="utf-8")
            try:
                app.write_text("{not-json", encoding="utf-8")
                loaded = load_templates(app)
                self.assertEqual(loaded, [])
                csv = self._write_csv("ok.csv", [["name"], ["x"]])
                result = run_worker([csv], self.output, workers=1)
                self.assertGreaterEqual(result["recordsExtracted"], 1)
            finally:
                if backup is None:
                    try:
                        app.unlink()
                    except OSError:
                        pass
                else:
                    app.write_text(backup, encoding="utf-8")
        finally:
            if previous_local_app_data is None:
                os.environ.pop("LOCALAPPDATA", None)
            else:
                os.environ["LOCALAPPDATA"] = previous_local_app_data

    def test_txt_rtf_xlsx_docx_pptx_csv(self):
        (self.temp / "note.txt").write_text("Invoice INV-77\nClient: Test Co\nTotal: 10.00", encoding="utf-8")
        (self.temp / "note.rtf").write_text(r"{\rtf1\ansi Invoice INV-88}", encoding="utf-8")
        book = Workbook()
        book.active.append(["name", "qty"])
        book.active.append(["widget", "3"])
        xlsx = self.temp / "grid.xlsx"
        book.save(xlsx)
        xlsm = self.temp / "macro.xlsm"
        book.save(xlsm)
        book.close()
        from docx import Document
        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Item"
        table.cell(0, 1).text = "Qty"
        table.cell(1, 0).text = "Pen"
        table.cell(1, 1).text = "4"
        docx = self.temp / "table.docx"
        doc.save(docx)
        from pptx import Presentation
        from pptx.util import Inches
        pres = Presentation()
        slide = pres.slides.add_slide(pres.slide_layouts[6])
        shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(0.5), Inches(6), Inches(2)).table
        shape.cell(0, 0).text = "Name"
        shape.cell(0, 1).text = "City"
        shape.cell(1, 0).text = "Ali"
        shape.cell(1, 1).text = "Riyadh"
        pptx = self.temp / "table.pptx"
        pres.save(pptx)
        result = run_worker(
            [self.temp / "note.txt", self.temp / "note.rtf", xlsx, xlsm, docx, pptx],
            self.output, workers=2,
        )
        self.assertEqual(result["filesProcessed"], 6)
        self.assertEqual(result["filesFailed"], 0)
        self.assertGreaterEqual(result["recordsExtracted"], 4)

    def test_legacy_office_without_file_fails_cleanly(self):
        doc = self.temp / "old.doc"
        doc.write_bytes(b"OLE fake doc")
        csv = self._write_csv("ok.csv", [["h"], ["1"]])
        result = run_worker([doc, csv], self.output, workers=1)
        self.assertEqual(result["filesProcessed"], 2)
        self.assertGreaterEqual(result["recordsExtracted"], 1)

    def test_images_bmp_webp_jpg_png_tiff(self):
        image = Image.new("RGB", (400, 120), "white")
        draw = ImageDraw.Draw(image)
        draw.text((20, 40), "Invoice INV-55 Total 9.00", fill="black", font=_font(22))
        paths = []
        for name, fmt in [("a.png", "PNG"), ("b.jpg", "JPEG"), ("c.bmp", "BMP"), ("d.tif", "TIFF"), ("e.webp", "WEBP")]:
            path = self.temp / name
            image.save(path, fmt)
            paths.append(path)
        result = run_worker(paths, self.output, workers=2)
        self.assertEqual(result["filesProcessed"], 5)
        self.assertLess(result["filesFailed"], 5)

    def test_text_and_scanned_pdf(self):
        text_pdf = self.temp / "text.pdf"
        image = Image.new("RGB", (500, 200), "white")
        draw = ImageDraw.Draw(image)
        draw.text((30, 80), "Tax Invoice INV-12", fill="black", font=_font(28))
        image.save(text_pdf, "PDF")
        scan = self.temp / "scan.pdf"
        page2 = Image.new("RGB", (500, 200), "white")
        ImageDraw.Draw(page2).text((30, 80), "Page Two Item 1", fill="black", font=_font(24))
        image.save(scan, "PDF", save_all=True, append_images=[page2])
        result = run_worker([text_pdf, scan], self.output, workers=1)
        self.assertEqual(result["filesProcessed"], 2)

    def test_batch_100_files(self):
        files = [self._write_csv(f"f{index:03}.csv", [["n", "q"], [str(index), "1"]]) for index in range(100)]
        result = run_worker(files, self.output, workers=min(4, os.cpu_count() or 2))
        self.assertEqual(result["filesProcessed"], 100)
        self.assertEqual(result["filesFailed"], 0)
        self.assertEqual(result["recordsExtracted"], 100)
        leftovers = list(self.output.glob(".*.sqlite"))
        self.assertEqual(leftovers, [])


class InvoiceExportTests(unittest.TestCase):
    def test_three_sheets_and_sku(self):
        rows = [
            ["Tax Invoice", "", "", "", ""],
            ["Customer: Armaan", "Invoice No: 118", "Date: 12-02-2020", "TRN: 121111111111111", ""],
            ["Description", "SKU", "Qty", "Unit Price", "Total"],
            ["Fridge", "FR-1", "2", "10.00", "20.00"],
            ["Subtotal: 20.00", "", "", "", ""],
        ]
        parsed = parse_invoice_table(rows)
        self.assertEqual(len(parsed["items"]), 1)
        self.assertEqual(parsed["items"][0]["sku"], "FR-1")
        self.assertEqual(parsed["header"].get("invoice_number"), "118")
        from export import write_invoice
        temp = Path(tempfile.mkdtemp(prefix="vertex-inv-"))
        dest = temp / "invoice.xlsx"
        try:
            write_invoice(dest, Path("sample.pdf"), parsed)
            book = load_workbook(dest)
            self.assertEqual(book.sheetnames, ["Header", "Items", "Totals"])
            self.assertEqual(book["Items"]["B1"].value, "SKU")
            self.assertEqual(book["Items"]["A2"].value, "Fridge")
            self.assertEqual(book["Totals"]["A2"].value, "Subtotal")
            self.assertEqual(book["Totals"]["A4"].value, "Grand Total")
            book.close()
        finally:
            shutil.rmtree(temp, ignore_errors=True)


class TableDetectionTests(unittest.TestCase):
    def test_projection_grid_has_regular_columns(self):
        import cv2
        import numpy as np
        from table_detect import detect_cells_from_projections
        image = np.full((240, 480), 255, np.uint8)
        for x in range(0, 481, 80):
            cv2.line(image, (x, 0), (x, 239), 0, 1)
        for y in range(0, 241, 40):
            cv2.line(image, (0, y), (479, y), 0, 1)
        rows = detect_cells_from_projections(image)
        self.assertGreaterEqual(len(rows), 4)
        self.assertGreaterEqual(len(rows[0]), 4)
        self.assertTrue(all(len(row) == len(rows[0]) for row in rows))

    def test_projection_uneven_columns_are_not_equal_splits(self):
        import cv2
        import numpy as np
        from table_detect import detect_cells_from_projections
        image = np.full((200, 500), 255, np.uint8)
        xs = [0, 50, 220, 310, 499]
        ys = [0, 40, 80, 120, 160, 199]
        for x in xs:
            cv2.line(image, (x, 0), (x, 199), 0, 2)
        for y in ys:
            cv2.line(image, (0, y), (499, y), 0, 2)
        rows = detect_cells_from_projections(image)
        self.assertGreaterEqual(len(rows), 3)
        self.assertEqual(len(rows[0]), 4)
        widths = [box[2] for box in rows[0]]
        self.assertGreater(max(widths) / max(min(widths), 1), 1.6)

    def test_language_follows_this_page_not_a_prior_layout(self):
        from ocr import choose_ocr_lang
        english_words = [
            {"text": "STATE", "conf": 90},
            {"text": "ALABAMA", "conf": 92},
            {"text": "1970", "conf": 88},
            {"text": "Capital", "conf": 91},
        ]
        self.assertEqual(choose_ocr_lang("", words=english_words * 8), "eng")
        arabic_words = [
            {"text": "الاسم", "conf": 90},
            {"text": "العنوان", "conf": 88},
            {"text": "الهاتف", "conf": 91},
            {"text": "شركة", "conf": 87},
        ]
        self.assertEqual(choose_ocr_lang("", words=arabic_words * 8), "ara+eng")
        noisy = [{"text": "خب", "conf": 12}, {"text": "STATE", "conf": 95}, {"text": "WATER", "conf": 93}]
        self.assertEqual(choose_ocr_lang("", words=noisy * 12), "eng")

    def test_numeric_ocr_tidying_and_decimals(self):
        from ocr import _choose_cell_text, _restore_column_decimals, _tidy_numeric_cell
        self.assertEqual(_tidy_numeric_cell("= 18032.67"), "18032.67")
        self.assertEqual(_tidy_numeric_cell("$.2"), "5.2")
        self.assertEqual(_tidy_numeric_cell("37299.9]"), "37299.91")
        self.assertEqual(_tidy_numeric_cell("1$032.67"), "15032.67")
        self.assertEqual(_tidy_numeric_cell("19,636,781"), "19636781")
        from ocr import _looks_numeric
        self.assertFalse(_looks_numeric("AED 100.00"))
        self.assertFalse(_looks_numeric("9/30/2020"))
        self.assertTrue(_looks_numeric("15032.67"))
        table = [
            ["STATE", "P_CAP", "UNEMP"],
            ["ALABAMA", "15032.67", "4.7"],
            ["ALABAMA", "15501.94", "5.2"],
            ["ALABAMA", "15972.41", "4.7"],
            ["ALABAMA", "16406.26", "3.9"],
            ["ALABAMA", "1676267", "74"],
            ["ALABAMA", "1912200", "11.0"],
            ["ARIZONA", "1159826", "42"],
        ]
        restored = _restore_column_decimals(table)
        self.assertEqual(restored[5][1], "16762.67")
        self.assertEqual(restored[6][1], "19122.00")
        self.assertEqual(restored[7][1], "11598.26")
        self.assertEqual(restored[5][2], "7.4")
        self.assertEqual(restored[7][2], "4.2")
        self.assertEqual(_choose_cell_text("= 18032.67", "1$032.67", 2), "1$032.67")
        self.assertEqual(_choose_cell_text("1925747", "9257.47", 2), "1925747")
        from ocr import _fix_isolated_letter_digits
        indexed = _fix_isolated_letter_digits([
            ["#", "STATE"],
            ["1", "ALABAMA"],
            ["2", "ALABAMA"],
            ["S", "ALABAMA"],
            ["4", "ALABAMA"],
            ["6", "ALABAMA"],
        ])
        self.assertEqual(indexed[3][0], "5")

    def test_word_clustering_does_not_mix_columns(self):
        words = [
            {"text": "Name", "left": 10, "top": 10, "width": 40, "height": 12, "conf": 90},
            {"text": "Qty", "left": 200, "top": 10, "width": 30, "height": 12, "conf": 90},
            {"text": "Ali", "left": 12, "top": 40, "width": 30, "height": 12, "conf": 88},
            {"text": "3", "left": 205, "top": 40, "width": 12, "height": 12, "conf": 92},
        ]
        table = cluster_words_to_table(words)
        self.assertEqual(table[0][0], "Name")
        self.assertEqual(table[0][1], "Qty")
        self.assertEqual(table[1][0], "Ali")
        self.assertEqual(table[1][1], "3")


class TemplateTests(unittest.TestCase):
    def test_save_reload_and_match(self):
        temp = Path(tempfile.mkdtemp(prefix="vertex-tpl-"))
        try:
            from templates import fingerprint
            header = "Gulf Enterprises Tax Invoice"
            save_templates(temp, [{"source": "a.pdf", "type": "invoice", "fingerprint": fingerprint(header), "header_text": header, "columns": ["Description"]}])
            loaded = load_templates(temp / "templates.json", app_templates_path())
            match = match_template(header, loaded)
            self.assertIsNotNone(match)
            self.assertEqual(match["type"], "invoice")
            self.assertIsNone(match_template("Contact directory email phone website", loaded))
        finally:
            shutil.rmtree(temp, ignore_errors=True)

    def test_classify_invoice_vs_table(self):
        self.assertEqual(classify_document("فاتورة ضريبية رقم الفاتورة 118", 5, 8), "invoice")
        self.assertEqual(classify_document("name city region", 4, 20), "table")
        # Short token "trn" must not hijack a contact/directory spreadsheet.
        contact = (
            "الاسم عنوان الايميل رقم الهاتف المحمول العنوان موقع الكتروني Instagram\n"
            "شركة ABC info@example.sa 966512345678 Riyadh Saudi Arabia https://example.sa https://instagram.com/x\n"
            "شركة XYZ hr@wind.sa 966555326098 Madinah https://linktr.ee/x https://instagram.com/y\n"
            "central contracting trn-noise-should-not-matter"
        )
        self.assertEqual(classify_document(contact, 6, 30), "table")

    def test_empty_cell_preserved_in_cluster(self):
        words = [
            {"text": "Name", "left": 10, "top": 10, "width": 40, "height": 12, "conf": 90},
            {"text": "Email", "left": 200, "top": 10, "width": 40, "height": 12, "conf": 90},
            {"text": "Phone", "left": 400, "top": 10, "width": 40, "height": 12, "conf": 90},
            {"text": "Ali", "left": 12, "top": 40, "width": 30, "height": 12, "conf": 88},
            {"text": "9665", "left": 405, "top": 40, "width": 40, "height": 12, "conf": 92},
        ]
        table = cluster_words_to_table(words)
        self.assertEqual(len(table[0]), 3)
        self.assertEqual(table[1][0], "Ali")
        self.assertEqual(table[1][1], "")
        self.assertEqual(table[1][2], "9665")

    def test_normalize_grid_forces_same_column_count(self):
        from table_detect import normalize_cell_grid
        rows = [
            [(0, 0, 40, 20), (50, 0, 40, 20), (100, 0, 40, 20)],
            [(0, 30, 40, 20), (50, 30, 40, 20)],
            [(0, 60, 40, 20), (50, 60, 40, 20), (100, 60, 40, 20), (150, 60, 40, 20)],
            [(0, 90, 40, 20), (50, 90, 40, 20), (100, 90, 40, 20)],
        ]
        normalized = normalize_cell_grid(rows, 200)
        self.assertTrue(all(len(row) == 3 for row in normalized))


class ReviewTests(unittest.TestCase):
    def test_apply_review_writes_excel(self):
        temp = Path(tempfile.mkdtemp(prefix="vertex-rev-"))
        xlsx = temp / "out.xlsx"
        book = Workbook()
        book.active.append(["Description"])
        book.active.append(["OLD"])
        book.save(xlsx)
        book.close()
        queue = temp / "review.json"
        queue.write_text(json.dumps({
            "items": [{"output": str(xlsx), "sheet": "Sheet", "row": 2, "column": 1, "value": "NEW"}]
        }), encoding="utf-8")
        env = os.environ.copy()
        env["PYTHONPATH"] = str(ROOT)
        completed = subprocess.run(
            [sys.executable, str(ROOT / "main.py"), "--apply-review", str(queue)],
            capture_output=True, text=True, env=env, timeout=60,
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        book = load_workbook(xlsx)
        self.assertEqual(book.active["A2"].value, "NEW")
        book.close()
        shutil.rmtree(temp, ignore_errors=True)


class RegressionTests(unittest.TestCase):
    def test_complex_and_comprehensive_expected_results(self):
        testdata = ROOT / "testdata"
        output = Path(tempfile.mkdtemp(prefix="vertex-reg-"))
        try:
            result = run_worker(
                [testdata / "Complex_Table_Report.xlsx", testdata / "Comprehensive_OCR_Test_Data.xlsx"],
                output, workers=2,
            )
            self.assertEqual(result["filesProcessed"], 2)
            self.assertEqual(result["filesFailed"], 0)
            self.assertEqual(result["recordsExtracted"], 5500)
            complex_out = next(output.glob("Complex_Table_Report_cleaned_*.xlsx"))
            book = load_workbook(complex_out, read_only=True, data_only=True)
            rows = list(book.active.iter_rows(values_only=True))
            book.close()
            headers = [str(value or "") for value in rows[0]]
            first = [str(value or "") for value in rows[1]]
            second = [str(value or "") for value in rows[2]]
            self.assertEqual(headers[1], "رمز الشركة (ID)")
            self.assertEqual(first[1], "CMP-95-00")
            self.assertTrue("طن" in first[3] or "tons" in second[3])
            comprehensive = next(output.glob("Comprehensive_OCR_Test_Data_cleaned_*.xlsx"))
            book = load_workbook(comprehensive, read_only=True, data_only=True)
            row = list(book.active.iter_rows(values_only=True))[1]
            book.close()
            self.assertEqual(str(row[1]), "046620")
        finally:
            shutil.rmtree(output, ignore_errors=True)


if __name__ == "__main__":
    unittest.main()
