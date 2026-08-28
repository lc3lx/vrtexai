"""Route each file extension to structured extraction or OCR pages."""
from __future__ import annotations

import io
import shutil
import tempfile
from pathlib import Path
from typing import Any

import numpy as np

from common import IMAGE_TYPES, OFFICE_TYPES, TABULAR_TYPES, TEXT_TYPES
from invoice import lines_to_table
from ocr import pdf_text_pages


def _ps_literal(path: str) -> str:
    return path.replace("'", "''")


def office_available(prog_id: str) -> bool:
    import subprocess
    script = (
        f"try {{ $null = New-Object -ComObject {prog_id}; exit 0 }} catch {{ exit 1 }}"
    )
    try:
        completed = subprocess.run(
            ["powershell", "-NoProfile", "-Command", script],
            capture_output=True, text=True, timeout=20,
        )
    except (OSError, subprocess.TimeoutExpired):
        return False
    return completed.returncode == 0


def convert_with_office(source: Path, dest_suffix: str) -> Path:
    import subprocess
    prog_id = "Word.Application" if dest_suffix == ".docx" else "PowerPoint.Application"
    if not office_available(prog_id):
        raise RuntimeError(
            f"ملفات {source.suffix.upper()} القديمة تتطلب Microsoft Office. "
            "ملفات DOCX وPPTX تعمل بدون Office."
        )
    temp_dir = Path(tempfile.mkdtemp(prefix="vertex-office-"))
    dest = temp_dir / (source.stem + dest_suffix)
    src = _ps_literal(str(source.resolve()))
    out = _ps_literal(str(dest))
    if dest_suffix == ".docx":
        script = (
            f"$app = New-Object -ComObject Word.Application; $app.Visible = $false; "
            f"$doc = $app.Documents.Open('{src}'); $doc.SaveAs([ref]'{out}', [ref]16); "
            f"$doc.Close(); $app.Quit()"
        )
    else:
        script = (
            f"$app = New-Object -ComObject PowerPoint.Application; "
            f"$pres = $app.Presentations.Open('{src}', $true, $false, $false); "
            f"$pres.SaveAs('{out}', 24); $pres.Close(); $app.Quit()"
        )
    try:
        completed = subprocess.run(
            ["powershell", "-NoProfile", "-Command", script],
            capture_output=True, text=True, timeout=120,
        )
    except subprocess.TimeoutExpired as error:
        shutil.rmtree(temp_dir, ignore_errors=True)
        raise RuntimeError(f"انتهت مهلة تحويل {source.suffix} عبر Microsoft Office.") from error
    if completed.returncode != 0 or not dest.is_file():
        shutil.rmtree(temp_dir, ignore_errors=True)
        detail = (completed.stderr or completed.stdout or "Office COM failed").strip()
        raise RuntimeError(f"تعذر تحويل {source.suffix} عبر Microsoft Office: {detail[:300]}")
    return dest


def extract_docx_tables(path: Path) -> tuple[list[list[list[str]]], list[str], list[np.ndarray]]:
    from docx import Document
    from PIL import Image
    document = Document(str(path))
    tables = [[[cell.text.strip() for cell in row.cells] for row in table.rows] for table in document.tables]
    paragraphs = [para.text.strip() for para in document.paragraphs if para.text.strip()]
    images: list[np.ndarray] = []
    import zipfile
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist():
            if not name.startswith("word/media/"):
                continue
            try:
                with Image.open(io.BytesIO(archive.read(name))) as image:
                    images.append(np.array(image.convert("RGB")))
            except Exception:
                continue
    return tables, paragraphs, images


def extract_pptx_tables(path: Path) -> tuple[list[list[list[str]]], list[str], list[np.ndarray]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    presentation = Presentation(str(path))
    tables: list[list[list[str]]] = []
    paragraphs: list[str] = []
    images: list[np.ndarray] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    paragraphs.append(text)
            if shape.has_table:
                table = shape.table
                tables.append([[cell.text.strip() for cell in row.cells] for row in table.rows])
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    with Image.open(io.BytesIO(blob)) as image:
                        images.append(np.array(image.convert("RGB")))
                except Exception:
                    continue
    return tables, paragraphs, images


def extract_text_file(path: Path) -> list[str]:
    raw = path.read_bytes()
    text = ""
    if path.suffix.lower() == ".rtf":
        try:
            from striprtf.striprtf import rtf_to_text
            text = rtf_to_text(raw.decode("utf-8", errors="ignore"))
        except Exception:
            text = raw.decode("utf-8", errors="ignore")
    else:
        for encoding in ("utf-8-sig", "utf-8", "cp1256", "cp1252"):
            try:
                text = raw.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            text = raw.decode("utf-8", errors="ignore")
    return [text]


def ingest(source: Path) -> dict[str, Any]:
    suffix = source.suffix.lower()
    payload: dict[str, Any] = {"tables": [], "paragraphs": [], "images": [], "warnings": []}
    if suffix in TABULAR_TYPES:
        payload["kind"] = "tabular"
        return payload
    if suffix == ".pdf":
        has_text, pages = pdf_text_pages(source)
        if has_text:
            payload["kind"] = "pdf_text"
            payload["tables"] = [lines_to_table(pages)]
            payload["paragraphs"] = pages
            return payload
        payload["kind"] = "images"
        payload["stream_images"] = True
        return payload
    if suffix in IMAGE_TYPES:
        payload["kind"] = "images"
        payload["stream_images"] = True
        return payload
    if suffix in TEXT_TYPES:
        pages = extract_text_file(source)
        payload["kind"] = "text"
        payload["tables"] = [lines_to_table(pages)]
        payload["paragraphs"] = pages
        return payload
    if suffix in OFFICE_TYPES:
        working = source
        converted_dir = None
        try:
            if suffix == ".doc":
                working = convert_with_office(source, ".docx")
                converted_dir = working.parent
            elif suffix == ".ppt":
                working = convert_with_office(source, ".pptx")
                converted_dir = working.parent
            if working.suffix.lower() == ".docx":
                tables, paragraphs, images = extract_docx_tables(working)
            else:
                tables, paragraphs, images = extract_pptx_tables(working)
        finally:
            if converted_dir is not None:
                shutil.rmtree(converted_dir, ignore_errors=True)
        payload["tables"] = tables
        payload["paragraphs"] = paragraphs
        payload["images"] = images
        payload["kind"] = "office"
        if converted_dir is not None:
            payload["warnings"].append("تم تحويل الملف القديم عبر Microsoft Office.")
        return payload
    payload["kind"] = "unsupported"
    return payload
